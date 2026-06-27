#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Générateur de Rapport de Fin d'Études en PowerPoint
Projet FlightHotel - Plateforme de Réservation Unifiée
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from datetime import datetime

# Créer la présentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Palette de couleurs
COLOR_PRIMARY = RGBColor(0, 120, 215)  # Bleu
COLOR_SECONDARY = RGBColor(68, 180, 57)  # Vert
COLOR_ACCENT = RGBColor(240, 100, 50)  # Orange
COLOR_DARK = RGBColor(30, 30, 30)
COLOR_LIGHT = RGBColor(240, 240, 240)

def add_title_slide(title, subtitle, author="", date=""):
    """Ajouter une slide de titre"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Sous-titre
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.1), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(200, 220, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Auteur et date
    if author or date:
        info_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.8))
        info_frame = info_box.text_frame
        p = info_frame.paragraphs[0]
        p.text = f"{author} | {date}" if author and date else (author or date)
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(200, 220, 255)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(title, content_points):
    """Ajouter une slide de contenu standard"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # En-tête coloré
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.color.rgb = COLOR_PRIMARY
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Contenu
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_DARK
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        p.level = 0
    
    return slide

def add_two_column_slide(title, left_title, left_content, right_title, right_content):
    """Ajouter une slide à deux colonnes"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # En-tête
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.color.rgb = COLOR_PRIMARY
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Colonne gauche
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(4.5), Inches(6.2))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    p = left_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    for line in left_content:
        p = left_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(15)
        p.space_before = Pt(6)
        p.level = 0
    
    # Colonne droite
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.0), Inches(4.5), Inches(6.2))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    p = right_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY
    
    for line in right_content:
        p = right_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(15)
        p.space_before = Pt(6)
        p.level = 0
    
    return slide

# ============================================
# SLIDE 1: Page de Couverture
# ============================================
add_title_slide(
    "RAPPORT DE FIN D'ÉTUDES",
    "FlightHotel\nPlateforme de Réservation Unifiée de Vols & d'Hôtels",
    "Auteur: Yacine Drissi",
    f"Date: {datetime.now().strftime('%d/%m/%Y')}"
)

# ============================================
# SLIDE 2: Table des Matières
# ============================================
add_content_slide("Table des Matières", [
    "① Introduction & Contexte",
    "② Architecture Globale de la Plateforme",
    "③ Stack Technique - Frontend",
    "④ Stack Technique - Backend",
    "⑤ Modélisation de la Base de Données",
    "⑥ Système de Contrôle d'Accès (RBAC)",
    "⑦ Fonctionnalités Clés",
    "⑨ Sécurité & Authentification",
    "⑩ Déploiement & Infrastructure",
    "⑪ Conclusion & Évolutions Futures"
])

# ============================================
# SLIDE 3: Introduction & Contexte
# ============================================
add_content_slide("Introduction & Contexte", [
    "🎯 Problématique:",
    "  • Les utilisateurs naviguent entre plusieurs plateformes pour réserver vols et hôtels",
    "  • Manque d'intégration unifiée dans l'écosystème du tourisme numérique",
    "",
    "💡 Solution: FlightHotel",
    "  • SaaS unifiée combinant recherche/réservation de vols & hôtels",
    "  • Interface fluide et moderne avec support multi-langue (FR/EN/AR)",
    "  • Gestion multi-rôles (Clients, Managers, Admins)"
])

# ============================================
# SLIDE 4: Architecture Globale
# ============================================
add_two_column_slide(
    "Architecture Globale",
    "Frontend",
    [
        "• Next.js 14+ App Router",
        "• Tailwind CSS (Responsive)",
        "• Zustand (State Management)",
        "• React Query v5 (Data Sync)",
        "• Support RTL (Arabe)"
    ],
    "Backend & Données",
    [
        "• Node.js + Express.js",
        "• TypeScript (Type-safe)",
        "• Prisma ORM",
        "• MySQL Database",
        "• JWT Authentication"
    ]
)

# ============================================
# SLIDE 5: Frontend - Next.js & Outils
# ============================================
add_content_slide("Frontend - Stack Technique", [
    "🔹 Next.js 14+ (App Router)",
    "  • Routage déclaratif avec dossiers dynamiques",
    "  • Optimisations de rendu côté client",
    "",
    "🔹 Tailwind CSS + Support RTL",
    "  • Responsive design avec directives RTL (rtl:text-right, rtl:flex-row-reverse)",
    "  • Passage dynamique du document en mode RTL pour l'arabe",
    "",
    "🔹 Zustand - Gestion d'État Légère",
    "  • Maintien du contexte utilisateur et tokens JWT",
    "",
    "🔹 React Query v5 - Synchronisation Async",
    "  • Cache automatique et rafraîchissements en arrière-plan"
])

# ============================================
# SLIDE 6: Backend - Node.js & Express
# ============================================
add_two_column_slide(
    "Backend - Architecture",
    "Express.js Framework",
    [
        "• Minimaliste et performant",
        "• Routage REST structuré",
        "• Middlewares CORS & Rate Limiting",
        "• Validation avec Express Validator",
        "• Gestion des erreurs globale"
    ],
    "TypeScript & Compilation",
    [
        "• ES2020 Target",
        "• CommonJS Modules",
        "• Strict Type Checking",
        "• Source Maps & Declarations",
        "• ts-node-dev pour le dev"
    ]
)

# ============================================
# SLIDE 7: Base de Données - Prisma ORM
# ============================================
add_content_slide("Base de Données - Prisma ORM", [
    "📊 Avantages de Prisma:",
    "  • Migrations automatisées & versionées",
    "  • Typage fort des requêtes au compile-time",
    "  • Réductions des bugs liés aux types",
    "",
    "🗄️ Schéma MySQL - Entités Principales:",
    "  • Users (avec Rôles: CLIENT, HOTEL_MANAGER, FLIGHT_MANAGER, ADMIN)",
    "  • AirlineCompany & Flights",
    "  • Hotel, Room, Reviews",
    "  • FlightReservation & HotelReservation",
    "  • Payment, Notification, PasswordReset"
])

# ============================================
# SLIDE 8: Entités & Relations
# ============================================
add_content_slide("Modélisation - Relations Clés", [
    "👤 User → FlightReservation/HotelReservation",
    "  Relation: One-to-Many (Un utilisateur, plusieurs réservations)",
    "",
    "✈️ Flight → AirlineCompany",
    "  Relation: Many-to-One (Plusieurs vols par compagnie)",
    "",
    "🏨 Hotel → Room → HotelReservation",
    "  Relation: One-to-Many (Cascades de suppressions)",
    "",
    "💳 Payment → FlightReservation / HotelReservation",
    "  Relation: One-to-One (Un paiement par réservation)",
    "",
    "⭐ Review → User & Flight",
    "  Constraint: @unique([userId, flightId])"
])

# ============================================
# SLIDE 9: Système RBAC
# ============================================
add_two_column_slide(
    "Système de Contrôle d'Accès (RBAC)",
    "Rôles Utilisateurs",
    [
        "👨‍💼 CLIENT",
        "  - Recherche & réservation",
        "  - Génération PDF",
        "  - Avis & évaluations",
        "",
        "🏨 HOTEL_MANAGER",
        "  - Gestion des chambres",
        "  - Dashboard hôtelier",
        "",
        "✈️ FLIGHT_MANAGER",
        "  - Gestion des vols",
        "  - Taux de remplissage",
        "",
        "🔐 ADMIN",
        "  - Contrôle complet"
    ],
    "Middlewares de Sécurité",
    [
        "✅ JWT Verification",
        "  - Tokens signés avec clés secrètes",
        "",
        "✅ Role-Based Authorization",
        "  - Vérification des rôles à chaque route",
        "",
        "✅ Rate Limiting",
        "  - 200 requêtes / 15 minutes",
        "",
        "✅ CORS Configuration",
        "  - Origin: localhost:3000 (dev)",
        "",
        "✅ Validation des Données",
        "  - Zod schemas (côté API)"
    ]
)

# ============================================
# SLIDE 10: Fonctionnalités Clés - Réservations
# ============================================
add_content_slide("Fonctionnalités - Système de Réservation", [
    "✈️ Réservation de Vols:",
    "  • Recherche multicritère (origine, destination, dates, passagers)",
    "  • Vérification de disponibilité des sièges",
    "  • Paiement simulé et confirmation",
    "",
    "🏨 Réservation d'Hôtels:",
    "  • Recherche par ville, dates et nombre de chambres",
    "  • Types de chambres (SINGLE, DOUBLE, SUITE, DELUXE)",
    "  • Gestion des prix et capacités",
    "",
    "📄 Génération de Documents:",
    "  • Billets d'avion (cartes d'embarquement)",
    "  • Bons de réservation (vouchers) avec QR Code",
    "  • PDFKit pour la création côté backend"
])

# ============================================
# SLIDE 11: Internationalisation & RTL
# ============================================
add_two_column_slide(
    "Internationalisation (i18n) & Support RTL",
    "Architecture légère",
    [
        "📌 React Context personnalisé",
        "  (LanguageProvider)",
        "",
        "📌 Pas de librairie tierce",
        "  lourde (ex: next-i18n)",
        "",
        "📌 Traduction récursive",
        "  t('common.navbar.home')",
        "",
        "📌 Substitution dynamique",
        "  de paramètres"
    ],
    "Support RTL (Arabe)",
    [
        "🌍 Langues supportées",
        "  • Français (FR)",
        "  • English (EN)",
        "  • العربية (AR)",
        "",
        "🔄 Direction dynamique",
        "  document.documentElement.dir = 'rtl'",
        "",
        "🎨 Tailwind RTL",
        "  rtl:text-right, rtl:flex-row-reverse",
        "",
        "📱 Selector responsive"
    ]
)

# ============================================
# SLIDE 12: Sécurité & Authentification
# ============================================
add_content_slide("Sécurité & Authentification", [
    "🔐 JWT (JSON Web Tokens):",
    "  • Tokens signés avec clés secrètes asymétriques",
    "  • Authentification stateless (sans état serveur)",
    "",
    "🛡️ Protections implémentées:",
    "  • Rate Limiting: 200 req/15 min par IP",
    "  • CORS: Origin validation stricte",
    "  • Express Validator: Sanitization des inputs",
    "  • Bcryptjs: Hash sécurisé des mots de passe",
    "",
    "📧 Récupération de Mot de Passe:",
    "  • Tokens UUIDv4 expirables",
    "  • Mailtrap + Nodemailer (SMTP sécurisé)",
    "  • Liens de réinitialisation avec TTL"
])

# ============================================
# SLIDE 13: Déploiement & Infrastructure
# ============================================
add_two_column_slide(
    "Déploiement & Infrastructure",
    "Docker & Compose",
    [
        "🐳 Containerization",
        "",
        "📦 Services Docker:",
        "  • MySQL 8.0 (pfe_mysql)",
        "  • phpMyAdmin (pfe_phpmyadmin)",
        "  • Backend Node.js",
        "",
        "💾 Volumes Persistants:",
        "  mysql_data:/"
        "var/lib/mysql"
    ],
    "Configuration & Variables",
    [
        "⚙️ Environment Variables:",
        "  • DATABASE_URL",
        "  • JWT_SECRET",
        "  • MAILTRAP_TOKEN",
        "  • SMTP_CONFIG",
        "",
        "🔧 Scripts NPM:",
        "  • npm run dev (Dev)",
        "  • npm run build (Prod)",
        "  • npm run prisma:migrate",
        "  • npm run prisma:seed"
    ]
)

# ============================================
# SLIDE 14: Module de Paiement
# ============================================
add_content_slide("Processus Transactionnel & Paiement", [
    "💳 Simulateur de Paiement:",
    "  • Prend en charge simulation de transactions par carte",
    "  • Stockage sécurisé du dernier chiffre (cardLast4)",
    "  • Statuts: PENDING → PAID / FAILED / REFUNDED",
    "",
    "📋 Génération PDF côté Backend:",
    "  • PDFKit intégré pour documents professionnels",
    "  • Polices modernes et formatage optimisé",
    "  • Inclusion de QR Codes de sécurité uniques",
    "",
    "✅ Flux de Paiement:",
    "  1. Création de Réservation (PENDING)",
    "  2. Initiation du Paiement",
    "  3. Confirmation & Génération de Documents",
    "  4. Historique & Téléchargement PDF"
])

# ============================================
# SLIDE 15: Dépendances Clés
# ============================================
add_content_slide("Dépendances Principales", [
    "Backend Dependencies:",
    "  • @prisma/client: ORM moderne",
    "  • express, cors: Framework web",
    "  • jsonwebtoken: JWT authentification",
    "  • bcryptjs: Hash sécurisé",
    "  • express-validator, zod: Validation",
    "  • nodemailer, mailtrap: Email SMTP",
    "  • pdfkit: Génération PDF",
    "  • express-rate-limit: Rate limiting",
    "",
    "DevDependencies:",
    "  • TypeScript 5.8: Typage statique",
    "  • ts-node-dev: Dev server avec hot-reload"
])

# ============================================
# SLIDE 16: Conclusion
# ============================================
add_content_slide("Conclusion", [
    "✅ Architecture moderne découpée (API-First)",
    "",
    "✅ Support multi-langue avec RTL complet",
    "",
    "✅ Système RBAC granulaire & sécurité renforcée",
    "",
    "✅ Flux transactionnels complexes gérés efficacement",
    "",
    "✅ Stack TypeScript pour maintenabilité & robustesse",
    "",
    "✅ Scalabilité via Prisma & structure modulaire"
])

# ============================================
# SLIDE 17: Évolutions Futures
# ============================================
add_content_slide("Perspectives d'Évolution", [
    "🔌 WebSockets:",
    "  • Synchronisation temps réel des sièges/chambres disponibles",
    "",
    "💳 Passerelle de Paiement Réelle:",
    "  • Intégration Stripe ou processeur bancaire en production",
    "",
    "🤖 Intelligence Artificielle:",
    "  • ML pour recommandations personnalisées d'hôtels",
    "  • Couplage intelligent avec destinations de vols",
    "",
    "📱 Applications Mobiles:",
    "  • React Native pour iOS/Android",
    "",
    "📊 Analytics & Monitoring:",
    "  • Dashboards temps réel pour stakeholders"
])

# ============================================
# SLIDE 18: Remerciements & Contact
# ============================================
add_title_slide(
    "Merci pour votre attention !",
    "Questions ?",
    "Yacine Drissi - FlightHotel Project",
    "GitHub: github.com/yacine20002"
)

# Sauvegarder la présentation
output_path = "Rapport_PFE_FlightHotel.pptx"
prs.save(output_path)
print(f"✅ Rapport généré avec succès: {output_path}")
print(f"📊 Total diapositives: {len(prs.slides)}")
